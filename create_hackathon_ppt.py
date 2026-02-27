from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_hackathon_ppt(filename="Behavioral_FIM_Pitch_Deck.pptx"):
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # 1. Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Behavioral FIM"
    subtitle.text = "AI-Powered File Integrity Monitoring\nHackathon Pitch Deck"
    
    # 2. The Problem
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "The Problem"
    tf = body_shape.text_frame
    tf.text = "Traditional FIM systems are noisy and rule-based."
    p = tf.add_paragraph()
    p.text = "Alert fatigue for SOC teams due to high false-positive rates."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Lack of context: A file change might be malicious or just a normal update."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Difficult to distinguish between benign administrative tasks and insider threats."
    p.level = 1

    # 3. The Solution
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Our Solution: Behavioral FIM"
    tf = body_shape.text_frame
    tf.text = "Context-aware file monitoring powered by Machine Learning."
    p = tf.add_paragraph()
    p.text = "Learns normal behavior for users and roles."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Highlights true anomalies instead of just 'any change'."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Interactive forensic dashboard for rapid incident response."
    p.level = 1
    
    # 4. Key Features
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Key Features"
    tf = body_shape.text_frame
    tf.text = "Real-time Monitoring & Alerting"
    p = tf.add_paragraph()
    p.text = "Dynamic Risk Scoring per user"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "ML Anomaly Detection (Isolation Forest)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Role-based Clustering (K-Means)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Correlation Engine for complex attack patterns"
    p.level = 0

    # 5. Architecture & Tech Stack
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Architecture & Tech Stack"
    tf = body_shape.text_frame
    tf.text = "Backend: Python, FastAPI"
    p = tf.add_paragraph()
    p.text = "Real-time updates: WebSockets"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Machine Learning: scikit-learn, numpy, pandas"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Frontend: Jinja2 Templates, HTML/CSS/JS (Glossy UI)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "System hooks: watchdog, psutil"
    p.level = 0

    # 6. Business Value & Impact
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Business Value & Impact"
    tf = body_shape.text_frame
    tf.text = "Drastically reduces false positives."
    p = tf.add_paragraph()
    p.text = "Speeds up threat hunting & incident response."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Protects critical infrastructure and ensures compliance."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Provides actionable intelligence, not just data logs."
    p.level = 1
    
    # 7. Future Roadmap / Next Steps
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Future Roadmap"
    tf = body_shape.text_frame
    tf.text = "Integration with SIEM solutions (Splunk, ELK)"
    p = tf.add_paragraph()
    p.text = "Advanced Deep Learning models for sequence prediction"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Automated remediation and rollback of malicious changes"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Cloud-native deployment (Kubernetes/Docker)"
    p.level = 0

    # 8. Q&A / Thank You
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Thank You!"
    subtitle.text = "Any Questions?\n\nDemo Time!"
    
    prs.save(filename)
    print(f"Presentation saved as {filename}")

if __name__ == "__main__":
    create_hackathon_ppt()
